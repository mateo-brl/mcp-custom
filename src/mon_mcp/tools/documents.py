"""
Module de generation de documents (Word, PowerPoint, PDF).

Necessite des dependances optionnelles:
- python-docx pour Word (.docx)
- python-pptx pour PowerPoint (.pptx)
- reportlab pour PDF
- markdown pour conversion Markdown

Outils: creer_word, creer_powerpoint, creer_pdf
"""

import json
import os
import re


def _parse_markdown_blocks(text: str) -> list[dict]:
    """Parse du Markdown simple en blocs structures."""
    blocks = []
    lines = text.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # Titres
        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            level = min(level, 6)
            blocks.append({"type": "heading", "level": level, "text": stripped.lstrip("# ").strip()})
            i += 1

        # Listes a puces
        elif stripped.startswith(("- ", "* ", "+ ")):
            items = []
            while i < len(lines) and lines[i].strip().startswith(("- ", "* ", "+ ")):
                items.append(lines[i].strip()[2:].strip())
                i += 1
            blocks.append({"type": "list", "items": items})

        # Listes numerotees
        elif re.match(r"^\d+\.\s", stripped):
            items = []
            while i < len(lines) and re.match(r"^\d+\.\s", lines[i].strip()):
                items.append(re.sub(r"^\d+\.\s*", "", lines[i].strip()))
                i += 1
            blocks.append({"type": "numbered_list", "items": items})

        # Paragraphe normal
        else:
            para_lines = []
            while i < len(lines) and lines[i].strip() and not lines[i].strip().startswith(("#", "- ", "* ", "+ ")):
                para_lines.append(lines[i].strip())
                i += 1
            blocks.append({"type": "paragraph", "text": " ".join(para_lines)})

    return blocks


def creer_word(chemin: str, contenu: str, formatage: str = "") -> str:
    """
    Cree un document Word (.docx) depuis du texte ou Markdown.

    Args:
        chemin: Chemin du fichier .docx a creer
        contenu: Contenu en texte ou Markdown (titres #, listes -, gras **, italique *)
        formatage: Options JSON optionnelles (police, taille, marges)

    Returns:
        Confirmation avec chemin, nombre de paragraphes, taille.
    """
    try:
        from docx import Document
        from docx.shared import Pt, Inches
    except ImportError:
        return json.dumps({
            "erreur": "python-docx non installe. Installez avec: pip install python-docx"
        }, ensure_ascii=False)

    chemin = os.path.abspath(chemin)
    if not chemin.endswith(".docx"):
        chemin += ".docx"

    # Options de formatage
    opts = {}
    if formatage:
        try:
            opts = json.loads(formatage)
        except json.JSONDecodeError:
            pass

    police = opts.get("police", "Calibri")
    taille = opts.get("taille", 11)

    try:
        doc = Document()

        # Style par defaut
        style = doc.styles["Normal"]
        font = style.font
        font.name = police
        font.size = Pt(taille)

        # Marges
        marges = opts.get("marges", {})
        if marges:
            section = doc.sections[0]
            if "top" in marges:
                section.top_margin = Inches(marges["top"])
            if "bottom" in marges:
                section.bottom_margin = Inches(marges["bottom"])
            if "left" in marges:
                section.left_margin = Inches(marges["left"])
            if "right" in marges:
                section.right_margin = Inches(marges["right"])

        # Parser le contenu Markdown
        blocks = _parse_markdown_blocks(contenu)
        para_count = 0

        for block in blocks:
            if block["type"] == "heading":
                level = min(block["level"], 4)
                doc.add_heading(block["text"], level=level)
                para_count += 1

            elif block["type"] == "list":
                for item in block["items"]:
                    doc.add_paragraph(item, style="List Bullet")
                    para_count += 1

            elif block["type"] == "numbered_list":
                for item in block["items"]:
                    doc.add_paragraph(item, style="List Number")
                    para_count += 1

            elif block["type"] == "paragraph":
                text = block["text"]
                p = doc.add_paragraph()
                # Gras et italique basiques
                parts = re.split(r"(\*\*.*?\*\*|\*.*?\*)", text)
                for part in parts:
                    if part.startswith("**") and part.endswith("**"):
                        run = p.add_run(part[2:-2])
                        run.bold = True
                    elif part.startswith("*") and part.endswith("*"):
                        run = p.add_run(part[1:-1])
                        run.italic = True
                    else:
                        p.add_run(part)
                para_count += 1

        os.makedirs(os.path.dirname(chemin), exist_ok=True)
        doc.save(chemin)

        taille_fichier = os.path.getsize(chemin)
        return json.dumps({
            "fichier": chemin,
            "paragraphes": para_count,
            "taille": f"{taille_fichier / 1024:.1f} KB" if taille_fichier > 1024 else f"{taille_fichier} B",
        }, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"erreur": str(e)}, ensure_ascii=False)


def creer_powerpoint(chemin: str, diapositives: str) -> str:
    """
    Cree une presentation PowerPoint (.pptx).

    Args:
        chemin: Chemin du fichier .pptx a creer
        diapositives: JSON - liste de diapos [{"titre": "...", "contenu": ["point1", "point2"], "notes": "..."}]

    Returns:
        Confirmation avec chemin, nombre de diapos, taille.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return json.dumps({
            "erreur": "python-pptx non installe. Installez avec: pip install python-pptx"
        }, ensure_ascii=False)

    chemin = os.path.abspath(chemin)
    if not chemin.endswith(".pptx"):
        chemin += ".pptx"

    try:
        slides_data = json.loads(diapositives) if isinstance(diapositives, str) else diapositives
        if not isinstance(slides_data, list):
            return json.dumps({"erreur": "Les diapositives doivent etre une liste JSON."}, ensure_ascii=False)
    except json.JSONDecodeError as e:
        return json.dumps({"erreur": f"JSON invalide: {e}"}, ensure_ascii=False)

    try:
        prs = Presentation()

        for slide_data in slides_data:
            titre = slide_data.get("titre", "")
            contenu = slide_data.get("contenu", [])
            notes = slide_data.get("notes", "")

            if contenu:
                # Layout avec titre + contenu
                layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = titre

                # Ajouter les points
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for i, point in enumerate(contenu):
                    if i == 0:
                        tf.text = str(point)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(point)
            else:
                # Layout titre seul
                layout = prs.slide_layouts[0]  # Title Slide
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = titre
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = ""

            # Notes du presentateur
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        os.makedirs(os.path.dirname(chemin), exist_ok=True)
        prs.save(chemin)

        taille_fichier = os.path.getsize(chemin)
        return json.dumps({
            "fichier": chemin,
            "diapositives": len(slides_data),
            "taille": f"{taille_fichier / 1024:.1f} KB" if taille_fichier > 1024 else f"{taille_fichier} B",
        }, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"erreur": str(e)}, ensure_ascii=False)


def creer_pdf(chemin: str, contenu: str, type_source: str = "markdown") -> str:
    """
    Cree un fichier PDF depuis du texte, Markdown ou HTML.

    Args:
        chemin: Chemin du fichier .pdf a creer
        contenu: Le contenu a convertir en PDF
        type_source: Type du contenu - "markdown", "text", ou "html" (defaut: "markdown")

    Returns:
        Confirmation avec chemin, nombre de pages, taille.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
    except ImportError:
        return json.dumps({
            "erreur": "reportlab non installe. Installez avec: pip install reportlab"
        }, ensure_ascii=False)

    chemin = os.path.abspath(chemin)
    if not chemin.endswith(".pdf"):
        chemin += ".pdf"

    try:
        os.makedirs(os.path.dirname(chemin), exist_ok=True)

        doc = SimpleDocTemplate(
            chemin,
            pagesize=A4,
            leftMargin=20 * mm,
            rightMargin=20 * mm,
            topMargin=20 * mm,
            bottomMargin=20 * mm,
        )

        styles = getSampleStyleSheet()
        story = []

        if type_source == "text":
            # Texte brut
            for line in contenu.split("\n"):
                if line.strip():
                    story.append(Paragraph(line, styles["Normal"]))
                else:
                    story.append(Spacer(1, 6 * mm))

        elif type_source in ("markdown", "html"):
            # Markdown -> blocs structures
            if type_source == "markdown":
                blocks = _parse_markdown_blocks(contenu)
            else:
                # HTML basique: juste l'envoyer tel quel comme paragraphes
                blocks = [{"type": "paragraph", "text": contenu}]

            for block in blocks:
                if block["type"] == "heading":
                    style_name = f"Heading{min(block['level'], 3)}"
                    story.append(Paragraph(block["text"], styles.get(style_name, styles["Heading1"])))
                    story.append(Spacer(1, 2 * mm))

                elif block["type"] in ("list", "numbered_list"):
                    for item in block["items"]:
                        bullet = "&bull; " if block["type"] == "list" else ""
                        story.append(Paragraph(f"    {bullet}{item}", styles["Normal"]))
                    story.append(Spacer(1, 3 * mm))

                elif block["type"] == "paragraph":
                    text = block["text"]
                    # Convertir Markdown gras/italique en HTML pour reportlab
                    text = re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", text)
                    text = re.sub(r"\*(.*?)\*", r"<i>\1</i>", text)
                    story.append(Paragraph(text, styles["Normal"]))
                    story.append(Spacer(1, 3 * mm))

        if not story:
            story.append(Paragraph("(Document vide)", styles["Normal"]))

        doc.build(story)

        taille_fichier = os.path.getsize(chemin)
        return json.dumps({
            "fichier": chemin,
            "pages": doc.page,
            "taille": f"{taille_fichier / 1024:.1f} KB" if taille_fichier > 1024 else f"{taille_fichier} B",
        }, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"erreur": str(e)}, ensure_ascii=False)


def register_tools(mcp):
    mcp.add_tool(creer_word)
    mcp.add_tool(creer_powerpoint)
    mcp.add_tool(creer_pdf)
